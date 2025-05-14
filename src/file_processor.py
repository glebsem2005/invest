import logging
import os
from abc import ABC, abstractmethod
from typing import BinaryIO, Dict, Type

import docx
import openpyxl
import pptx
import PyPDF2
from aiogram import Bot
from aiogram.types import Document

logger = logging.getLogger('bot')


class FileExtractor(ABC):
    """Абстрактный класс для извлечения текста из файлов различных форматов."""

    @abstractmethod
    async def extract_text(self, file: BinaryIO) -> str:
        """Извлекает текст из файла."""
        ...


class PDFExtractor(FileExtractor):
    """Извлечение текста из PDF файлов."""

    async def extract_text(self, file: BinaryIO) -> str:
        try:
            reader = PyPDF2.PdfReader(file)
            text = ''

            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text += page.extract_text() + '\n\n'

            return text.strip()
        except Exception as e:
            logger.error(f'Ошибка при извлечении текста из PDF: {e}')
            raise ValueError(f'Не удалось извлечь текст из PDF: {e}')


class DocxExtractor(FileExtractor):
    """Извлечение текста из DOCX файлов."""

    async def extract_text(self, file: BinaryIO) -> str:
        try:
            doc = docx.Document(file)
            text = ''

            for paragraph in doc.paragraphs:
                text += paragraph.text + '\n'

            return text.strip()
        except Exception as e:
            logger.error(f'Ошибка при извлечении текста из DOCX: {e}')
            raise ValueError(f'Не удалось извлечь текст из DOCX: {e}')


class PPTXExtractor(FileExtractor):
    """Извлечение текста из PPTX файлов."""

    async def extract_text(self, file: BinaryIO) -> str:
        try:
            presentation = pptx.Presentation(file)
            text = ''

            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'
                text += '\n'

            return text.strip()
        except Exception as e:
            logger.error(f'Ошибка при извлечении текста из PPTX: {e}')
            raise ValueError(f'Не удалось извлечь текст из PPTX: {e}')


class TXTExtractor(FileExtractor):
    """Извлечение текста из TXT файлов."""

    async def extract_text(self, file: BinaryIO) -> str:
        try:
            content = file.read().decode('utf-8')
            return content.strip()
        except UnicodeDecodeError:
            try:
                file.seek(0)
                content = file.read().decode('cp1251')
                return content.strip()
            except Exception as e:
                logger.error(f'Ошибка при извлечении текста из TXT: {e}')
                raise ValueError(f'Не удалось извлечь текст из TXT: {e}')


class ExcelExtractor(FileExtractor):
    """Извлечение текста из Excel файлов."""

    @staticmethod
    def _workbook_to_text(wb) -> str:
        text = ''
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                text += row_text + '\n'
        return text.strip()

    async def extract_text(self, file: BinaryIO) -> str:
        try:
            file.seek(0)
            wb = openpyxl.load_workbook(file, read_only=True, data_only=True)
            return self._workbook_to_text(wb)
        except Exception as e:
            logger.error(f'Ошибка при извлечении текста из Excel: {e}')
            raise ValueError(f'Не удалось извлечь текст из Excel: {e}')

    @staticmethod
    async def extract_text_from_path(path: str) -> str:
        try:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            return ExcelExtractor._workbook_to_text(wb)
        except Exception as e:
            logger.error(f'Ошибка при чтении Excel-файла по пути {path}: {e}')
            raise ValueError(f'Ошибка при чтении Excel-файла по пути {path}: {e}')


class FileProcessor:
    """Класс для обработки файлов различных форматов."""

    _extractors: Dict[str, Type[FileExtractor]] = {
        '.pdf': PDFExtractor,
        '.docx': DocxExtractor,
        '.doc': DocxExtractor,
        '.pptx': PPTXExtractor,
        '.ppt': PPTXExtractor,
        '.txt': TXTExtractor,
        '.xlsx': ExcelExtractor,
    }

    @classmethod
    async def extract_text_from_file(cls, document: Document, bot: Bot) -> str:
        """Извлекает текст из файла."""
        file_name = document.file_name.lower()
        file_ext = os.path.splitext(file_name)[1]

        if file_ext not in cls._extractors:
            supported_formats = ', '.join(cls._extractors.keys())
            raise ValueError(f'Формат файла {file_ext} не поддерживается. Поддерживаемые форматы: {supported_formats}')

        file_id = document.file_id
        file = await bot.get_file(file_id)
        file_path = file.file_path
        downloaded_file = await bot.download_file(file_path)

        extractor_class = cls._extractors[file_ext]
        extractor = extractor_class()

        try:
            text = await extractor.extract_text(downloaded_file)
            return text
        except Exception as e:
            logger.error(f'Ошибка при обработке файла {file_name}: {e}')
            raise ValueError(f'Ошибка при обработке файла: {e}')

    @classmethod
    def register_extractor(cls, extension: str, extractor: Type[FileExtractor]) -> None:
        """Регистрирует новый экстрактор для указанного расширения файла."""
        cls._extractors[extension.lower()] = extractor
